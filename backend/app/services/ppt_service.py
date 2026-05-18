"""PPT generation service: structured content → python-pptx → .pptx file.

Pipeline:
  topic + chunks → SlideDeck JSON (Mock or LLM)
  → python-pptx rendering
  → .pptx bytes
"""

import logging
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.schemas.resource import LectureDocJSON

logger = logging.getLogger(__name__)

# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette for academic slides
COLOR_TITLE = RGBColor(0x1A, 0x56, 0xDB)       # blue
COLOR_SUBTITLE = RGBColor(0x55, 0x55, 0x55)     # gray
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)         # dark gray
COLOR_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)       # light blue bg
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"

_PPT_PROMPT = """你是一位课件设计专家。请根据课程内容和学生画像，生成PPT课件大纲。
输出纯JSON（不要markdown标记），格式：
{{
  "title": "课件标题",
  "subtitle": "副标题",
  "slides": [
    {{
      "type": "title 或 content",
      "title": "本页标题",
      "bullets": ["要点1", "要点2"],
      "speaker_notes": "讲师备注"
    }}
  ]
}}
主题：{topic}
课程资料：
{context}
学生画像：
{profile}
要求：生成6-8页。包含封面、学习目标、前置知识、核心概念、例题、常见误区、练习建议、总结。内容基于课程资料。体现学生画像。只输出JSON。"""


def _clean_ppt_json(text: str) -> str:
    """Strip fences and fix common JSON issues."""
    import re
    text = text.strip()
    text = re.sub(r'^```(?:json)?\s*\n?', '', text)
    text = re.sub(r'\n?```\s*$', '', text)
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end > start:
        text = text[start:end+1]
    text = re.sub(r',\s*}', '}', text)
    text = re.sub(r',\s*]', ']', text)
    return text


def _chunks_to_context(chunks: list[dict], max_len: int = 800) -> str:
    """Merge chunks into prompt-ready context for PPT generation."""
    parts = []
    total = 0
    for i, c in enumerate(chunks):
        ct = (c.get("content", "") or "").strip()
        if not ct:
            continue
        src = c.get("source", "unknown")
        part = f"[{i+1}:{src}] {ct}"
        if total + len(part) > max_len:
            remaining = max_len - total
            parts.append(part[:remaining])
            break
        parts.append(part)
        total += len(part)
    return "\n\n".join(parts) if parts else "（暂无课程资料）"


def _build_profile_text(profile: dict) -> str:
    """Build profile summary for PPT prompts."""
    import json
    parts = []
    if profile.get("major"):
        parts.append(f"专业：{profile['major']}")
    parts.append(f"知识水平：{profile.get('knowledge_level', 'intermediate')}")
    parts.append(f"认知风格：{profile.get('cognitive_style', 'conceptual')}")
    parts.append(f"学习节奏：{profile.get('pace_preference', 'moderate')}")
    wps = profile.get("weak_points")
    if wps:
        if isinstance(wps, str):
            try: wps = json.loads(wps)
            except: pass
        if isinstance(wps, list) and wps:
            parts.append(f"知识薄弱点：{', '.join(wps)}")
    return "\n".join(parts)





# ── SlideDeck JSON generation ─────────────────────────────────────────────────

def build_slide_deck(
    topic: str,
    chunks: list[dict],
    student_profile: dict,
    llm_provider=None,
) -> dict:
    """Build SlideDeck JSON from topic and retrieved chunks.

    Uses DeepSeek LLM when available for slide content.
    Layout/template is fixed; slide text is LLM-generated.
    """
    # Try LLM for slide content
    if llm_provider and llm_provider.provider != "mock" and student_profile:
        try:
            context = _chunks_to_context(chunks, 800)
            profile_text = _build_profile_text(student_profile)
            prompt = _PPT_PROMPT.format(topic=topic, context=context, profile=profile_text)
            resp = llm_provider.generate(
                [{"role": "user", "content": prompt}], temperature=0.3
            )
            cleaned = _clean_ppt_json(resp.content.strip())
            import json as _json
            data = _json.loads(cleaned)
            slides = data.get("slides", [])
            if isinstance(slides, list) and len(slides) >= 4:
                result = {
                    "title": data.get("title", f"{topic}课件"),
                    "subtitle": data.get("subtitle", ""),
                    "slides": [
                        {
                            "type": s.get("type", "content"),
                            "title": s.get("title", ""),
                            "bullets": s.get("bullets", []),
                            "speaker_notes": s.get("speaker_notes", ""),
                        }
                        for s in slides
                        if isinstance(s, dict) and s.get("title")
                    ],
                    "generated_by": "deepseek",
                }
                if result["slides"]:
                    return result
        except Exception:
            pass

    # Fallback to template (Mock)
    summary = _extract_summary(chunks, 800)
    bullets_1 = _extract_bullets(summary, 4, 0)
    bullets_2 = _extract_bullets(summary, 4, 4)
    bullets_3 = _extract_bullets(summary, 4, 8) if len(summary) > 200 else _extract_bullets(summary, 4, 4)
    return {
        "title": f"{topic}",
        "subtitle": "智能学习Agent 自动生成",
        "slides": [
            {"type": "title", "title": f"{topic}", "subtitle": "智能学习Agent — 个性化课件", "bullets": [], "speaker_notes": ""},
            {"type": "content", "title": "学习目标", "bullets": bullets_1, "speaker_notes": "基于课程资料自动生成"},
            {"type": "content", "title": "核心概念", "bullets": bullets_2, "speaker_notes": "从上文提取的关键概念"},
            {"type": "content", "title": "重点与难点", "bullets": bullets_3, "speaker_notes": "需要额外练习的内容"},
            {"type": "content", "title": "练习建议", "bullets": ["完成课后习题", "制作思维导图", "做相关测验题", "复习易错概念"], "speaker_notes": ""},
            {"type": "content", "title": "总结", "bullets": ["回顾核心概念", "整理知识框架", "准备下一主题"], "speaker_notes": "感谢使用智能学习Agent"},
        ],
        "generated_by": "fallback_template",
    }
def _extract_summary(chunks: list[dict], max_len: int = 800) -> str:
    """Extract a combined summary from retrieved chunks."""
    parts = []
    total = 0
    for c in chunks:
        content = (c.get("content", "") or "").strip()
        if not content:
            continue
        parts.append(content)
        total += len(content)
        if total >= max_len:
            break
    return "\n".join(parts)


def _extract_bullets(text: str, count: int = 4, offset: int = 0) -> list[str]:
    """Extract bullet points from text by splitting on common delimiters."""
    import re
    # Split on Chinese periods, newlines, or numbered items
    sentences = re.split(r"[。；\n]|(?<=\d)[\.、)]", text)
    bullets = []
    for s in sentences:
        s = s.strip()
        if len(s) > 5 and len(s) < 80:
            bullets.append(s)
    start = min(offset, len(bullets))
    end = min(start + count, len(bullets))
    result = bullets[start:end]
    # Pad if not enough
    while len(result) < count:
        result.append("请参考课程教材获取更多内容")
    return result[:count]


# ── PPTX Rendering ────────────────────────────────────────────────────────────

def render_pptx(slide_deck: dict) -> bytes:
    """Render a SlideDeck JSON to .pptx file bytes using python-pptx.

    Args:
        slide_deck: SlideDeck JSON dict with title, subtitle, slides[].

    Returns:
        .pptx file bytes.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides = slide_deck.get("slides", [])
    for slide_data in slides:
        stype = slide_data.get("type", "content")
        if stype == "title":
            _add_title_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: Presentation, data: dict):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_TITLE

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.8), Inches(10), Inches(1.0)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xEE)
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, data: dict):
    """Add a content slide with title and bullets."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title bar
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.3),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_TITLE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.2), Inches(11), Inches(1.0)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_TITLE

    # Bullets
    bullets = data.get("bullets", [])
    if bullets:
        body_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(1.8), Inches(11), Inches(5.0)
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_BODY
            p.font.name = FONT_BODY
            p.space_after = Pt(12)
